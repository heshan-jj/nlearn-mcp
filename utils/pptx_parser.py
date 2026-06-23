import io
import logging

from pptx import Presentation

logger = logging.getLogger(__name__)


def _shape_text(shape) -> str:
    """Return all text from a shape, including table cells."""
    if shape.has_text_frame:
        lines = [para.text for para in shape.text_frame.paragraphs if para.text.strip()]
        return "\n".join(lines)

    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    return ""


def parse_pptx_content(pptx_bytes: bytes) -> str:
    """
    Extract text from a PPTX file, one block per slide separated by slide markers.
    """
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            shape_parts = [_shape_text(shape) for shape in slide.shapes]
            slide_content = "\n".join(p for p in shape_parts if p)
            if slide_content.strip():
                slides_text.append(f"--- Slide {i} ---\n{slide_content}")

        return "\n\n".join(slides_text)
    except Exception as e:
        logger.error("Error extracting PPTX: %s", e)
        return f"Error extracting text from PPTX: {e}"
