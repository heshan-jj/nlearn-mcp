import io
import unittest

from pptx import Presentation
from pptx.util import Inches

from utils.pptx_parser import parse_pptx_content


def _make_pptx(slides: list[list[str]]) -> bytes:
    """Create a minimal PPTX with one text box per string per slide."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    for slide_texts in slides:
        slide = prs.slides.add_slide(blank_layout)
        for text in slide_texts:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestParsePptxContent(unittest.TestCase):
    def test_extracts_slide_text(self) -> None:
        pptx_bytes = _make_pptx([["Title slide text"], ["Body slide text"]])
        result = parse_pptx_content(pptx_bytes)
        self.assertIn("Title slide text", result)
        self.assertIn("Body slide text", result)

    def test_includes_slide_markers(self) -> None:
        pptx_bytes = _make_pptx([["Slide A"], ["Slide B"]])
        result = parse_pptx_content(pptx_bytes)
        self.assertIn("--- Slide 1 ---", result)
        self.assertIn("--- Slide 2 ---", result)

    def test_skips_empty_slides(self) -> None:
        pptx_bytes = _make_pptx([[], ["Non-empty"]])
        result = parse_pptx_content(pptx_bytes)
        self.assertNotIn("--- Slide 1 ---", result)
        self.assertIn("--- Slide 2 ---", result)

    def test_handles_invalid_bytes(self) -> None:
        result = parse_pptx_content(b"not a pptx file")
        self.assertTrue(result.startswith("Error extracting text from PPTX:"))


if __name__ == "__main__":
    unittest.main()
